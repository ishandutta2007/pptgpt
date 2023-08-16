from pptx import Presentation

# Create a presentation object
presentation = Presentation()

# Slide 1 - Title slide
slide1 = presentation.slides.add_slide(presentation.slide_layouts[0])
title = slide1.shapes.title
title.text = "Improving the Performance of a Machine Learning Algorithm"

# Slide 2 - Introduction
slide2 = presentation.slides.add_slide(presentation.slide_layouts[1])
content2 = slide2.shapes.title
content2.text = "Introduction"

# Slide 3 - Balancing Techniques
slide3 = presentation.slides.add_slide(presentation.slide_layouts[1])
content3 = slide3.shapes.title
content3.text = "Balancing Techniques"

# Slide 4 - Using imbalanced-learn library
slide4 = presentation.slides.add_slide(presentation.slide_layouts[1])
content4 = slide4.shapes.title
content4.text = "Using imbalanced-learn Library"

# Slide 5 - Data Import
slide5 = presentation.slides.add_slide(presentation.slide_layouts[1])
content5 = slide5.shapes.title
content5.text = "Data Import"

# Slide 6 - Model Building
slide6 = presentation.slides.add_slide(presentation.slide_layouts[1])
content6 = slide6.shapes.title
content6.text = "Model Building"

# Slide 7 - Imbalanced Dataset
slide7 = presentation.slides.add_slide(presentation.slide_layouts[1])
content7 = slide7.shapes.title
content7.text = "Imbalanced Dataset"

# Set the content for each slide
slides = [slide2, slide3, slide4, slide5, slide6, slide7]
content = [
    "This tutorial belongs to the series How to improve the performance of a Machine Learning Algorithm.",
    "A balanced dataset is a dataset where each output class is represented by the same number of input samples.",
    "Balancing techniques include oversampling, undersampling, class weight, and threshold.",
    "The imbalanced-learn library, part of the contrib packages of scikit-learn, is used in this tutorial.",
    "Data is imported using the pandas library, and the target class is created based on cuisine.",
    "The model is built using the Decision Tree algorithm, and evaluation metrics are calculated.",
    "The classification report and various plots are generated for the imbalanced dataset."
]

# Add content to each slide
for slide, content_text in zip(slides, content):
    content_placeholder = slide.placeholders[1]
    content_placeholder.text = content_text

# Save the presentation
presentation.save("machine_learning_presentation.pptx")
